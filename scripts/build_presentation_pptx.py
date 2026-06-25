#!/usr/bin/env python3
"""
Build PPTX presentation for Airborne Gravity Data Processing project.
Output: output/presentation.pptx
"""

import os
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_DIR = os.path.normpath(os.path.join(SCRIPT_DIR, "..", "output"))
OUTPUT_FILE = os.path.join(OUTPUT_DIR, "presentation.pptx")

# ──────────────────────────────────────────────
# COLOR PALETTE
# ──────────────────────────────────────────────
NAVY      = RGBColor(0x1A, 0x2A, 0x5E)
TEAL      = RGBColor(0x02, 0x80, 0x90)
TEAL_MID  = RGBColor(0x0F, 0x9B, 0x8E)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF8, 0xFA, 0xFC)
LIGHT_BLU = RGBColor(0xEF, 0xF6, 0xFF)
DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
MUTED     = RGBColor(0x64, 0x74, 0x8B)
BORDER    = RGBColor(0xE2, 0xE8, 0xF0)
LIGHT_AQU = RGBColor(0xA8, 0xD8, 0xEA)

# ──────────────────────────────────────────────
# LAYOUT CONSTANTS  (inches)
# ──────────────────────────────────────────────
SW   = 10.0    # slide width
SH   = 5.625   # slide height
HDR  = 0.62    # header bar height
ACC  = 0.05    # teal accent line height
CTY  = HDR + ACC + 0.08   # content area top (≈ 0.75")
MARG = 0.35    # horizontal margin

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

# ──────────────────────────────────────────────
# PRIMITIVE HELPERS
# ──────────────────────────────────────────────

def set_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def rect(slide, x, y, w, h, fill=None, line_color=None, line_pt=0):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is not None:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_color is not None:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_pt) if line_pt else Pt(0.75)
    else:
        s.line.fill.background()
    return s


def vtxt(slide, x, y, w, h, text, size=14, bold=False, italic=False,
         color=None, align=PP_ALIGN.LEFT, font="Calibri", valign="ctr"):
    """Textbox with vertical alignment."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    bp = tf._txBody.find(qn("a:bodyPr"))
    if bp is not None:
        bp.set("anchor", valign)
        bp.set("wrap", "square")
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    if color:
        r.font.color.rgb = color
    return tb


def img_ratio(path):
    """Return (width, height, ratio) of image."""
    with PILImage.open(path) as im:
        iw, ih = im.size
    return iw, ih, iw / ih


def add_img(slide, path, x, y, max_w, max_h):
    """
    Add image top-left-aligned, constrained to (max_w × max_h), preserving ratio.
    Returns the bottom edge Y of the placed image.
    """
    iw, ih, ratio = img_ratio(path)
    if max_w / ratio <= max_h:
        w, h = max_w, max_w / ratio
    else:
        h, w = max_h, max_h * ratio
    # Horizontally center within max_w; top-align vertically
    cx = x + (max_w - w) / 2
    slide.shapes.add_picture(path, Inches(cx), Inches(y), Inches(w), Inches(h))
    return y + h   # bottom edge


def content_slide(title):
    """Add a slide with navy header bar + teal accent, return slide."""
    sl = prs.slides.add_slide(BLANK)
    set_bg(sl, OFF_WHITE)
    rect(sl, 0, 0, SW, HDR, fill=NAVY)
    rect(sl, 0, HDR, SW, ACC, fill=TEAL)
    vtxt(sl, MARG, 0, SW - MARG * 2, HDR, title,
         size=22, bold=True, color=WHITE, align=PP_ALIGN.LEFT, valign="ctr")
    return sl


def card(slide, x, y, w, h, accent_color=None):
    """White card with optional left accent bar."""
    rect(slide, x, y, w, h, fill=WHITE, line_color=BORDER, line_pt=0.75)
    if accent_color:
        rect(slide, x, y, 0.07, h, fill=accent_color)


def multi_para(slide, x, y, w, h, items, font="Calibri", default_size=13,
               default_color=None, default_bold=False, valign="t"):
    """
    items: list of dicts with keys: text, size, bold, italic, color,
           align (PP_ALIGN), space_after (Pt), new_para (bool).
    """
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    bp = tf._txBody.find(qn("a:bodyPr"))
    if bp is not None:
        bp.set("anchor", valign)
    first = True
    for item in items:
        t    = item.get("text", "")
        size = item.get("size",  default_size)
        bold = item.get("bold",  default_bold)
        it   = item.get("italic", False)
        col  = item.get("color", default_color or DARK_TEXT)
        aln  = item.get("align", PP_ALIGN.LEFT)
        sp   = item.get("space_after", 0)
        np   = item.get("new_para", True)
        if first:
            p = tf.paragraphs[0]
            first = False
        elif np:
            p = tf.add_paragraph()
        p.alignment = aln
        if sp:
            p.space_after = Pt(sp)
        r = p.add_run()
        r.text = t
        r.font.size   = Pt(size)
        r.font.bold   = bold
        r.font.italic = it
        r.font.name   = font
        r.font.color.rgb = col


# ──────────────────────────────────────────────
# SLIDE 1 — TITLE
# ──────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)
set_bg(s1, NAVY)
rect(s1, 0, 0,        SW, 0.18, fill=TEAL)
rect(s1, 0, SH - 0.18, SW, 0.18, fill=TEAL)

vtxt(s1, 0.5, 0.9, 9.0, 1.3,
     "Pengolahan Data Gravitasi Airborne",
     size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER, valign="ctr")

vtxt(s1, 0.5, 2.35, 9.0, 0.75,
     "Wilayah Sulawesi Tengah, Indonesia",
     size=22, color=LIGHT_AQU, align=PP_ALIGN.CENTER, valign="ctr")

rect(s1, 2.5, 3.2, 5.0, 0.04, fill=TEAL)

vtxt(s1, 0.5, 3.35, 9.0, 0.55,
     "[Student Name]  —  Progress Report",
     size=16, color=RGBColor(0xCB, 0xD5, 0xE1),
     align=PP_ALIGN.CENTER, valign="ctr")


# ──────────────────────────────────────────────
# SLIDE 2 — OUTLINE  (compact card height)
# ──────────────────────────────────────────────
s2 = content_slide("Outline")

outline = [
    "Pendahuluan & Data Survey",
    "Metodologi Pengolahan",
    "Quality Control Data",
    "Persiapan DEM (Digital Elevation Model)",
    "Anomali Bouguer",
    "Koreksi Terrain",
    "Gridding (Equivalent Sources)",
    "Peta Anomali Gravitasi",
    "Hasil & Interpretasi",
    "Kesimpulan & Rencana Selanjutnya",
]
half = (len(outline) + 1) // 2
# Cards fill the full content area; items spread with generous spacing
CARD_H = SH - CTY - 0.15   # fill slide height

for col_idx, (items, x_off, acc_col) in enumerate(
        [(outline[:half], 0.35, TEAL), (outline[half:], 5.25, NAVY)]):
    card(s2, x_off, CTY + 0.1, 4.5, CARD_H, accent_color=acc_col)
    start = 1 if col_idx == 0 else half + 1
    rows = [{"text": f"{start+i}.  {txt}", "size": 14,
             "color": DARK_TEXT, "space_after": 18}
            for i, txt in enumerate(items)]
    multi_para(s2, x_off + 0.18, CTY + 0.3, 4.2, CARD_H - 0.4, rows)


# ──────────────────────────────────────────────
# SLIDE 3 — PENDAHULUAN
# ──────────────────────────────────────────────
s3 = content_slide("Pendahuluan")

LEFT_CARD_H = SH - CTY - 0.2
card(s3, 0.35, CTY, 5.8, LEFT_CARD_H, accent_color=TEAL)
multi_para(s3, 0.55, CTY + 0.15, 5.5, LEFT_CARD_H - 0.3, [
    {"text": "Tujuan Penelitian",
     "size": 14, "bold": True, "color": NAVY, "space_after": 4},
    {"text": "Mengolah data gravitasi airborne untuk menghasilkan peta anomali "
             "gravitasi wilayah Sulawesi Tengah, Indonesia",
     "size": 12, "color": DARK_TEXT, "space_after": 12},
    {"text": "Wilayah Studi",
     "size": 14, "bold": True, "color": NAVY, "space_after": 4},
    {"text": "Sulawesi Tengah, Indonesia",
     "size": 12, "color": DARK_TEXT, "space_after": 2},
    {"text": "Longitude: 120.80 deg - 124.10 deg BT",
     "size": 12, "color": DARK_TEXT, "space_after": 2},
    {"text": "Latitude: 2.10 deg LS - 0.003 deg LU",
     "size": 12, "color": DARK_TEXT, "space_after": 12},
    {"text": "Data Pengamatan",
     "size": 14, "bold": True, "color": NAVY, "space_after": 4},
    {"text": "48,806 titik pengamatan gravitasi airborne",
     "size": 12, "color": DARK_TEXT, "space_after": 2},
    {"text": "Mencakup wilayah darat dan laut (mixed marine/land survey)",
     "size": 12, "color": DARK_TEXT},
])

# Right side: two stacked cards
INFO_CARD_H = 1.85
STAT_CARD_H = 1.8
CARD_GAP    = 0.15
rx = 6.45
rw = 3.2

card(s3, rx, CTY, rw, INFO_CARD_H, accent_color=TEAL_MID)
multi_para(s3, rx + 0.18, CTY + 0.12, rw - 0.25, INFO_CARD_H - 0.2, [
    {"text": "Lokasi Survei",
     "size": 12, "bold": True, "color": NAVY, "align": PP_ALIGN.CENTER, "space_after": 8},
    {"text": "Sulawesi Tengah",
     "size": 14, "bold": True, "color": TEAL, "align": PP_ALIGN.CENTER, "space_after": 5},
    {"text": "120.80 deg - 124.10 deg BT",
     "size": 11, "color": DARK_TEXT, "align": PP_ALIGN.CENTER, "space_after": 3},
    {"text": "2.10 deg LS - 0.003 deg LU",
     "size": 11, "color": DARK_TEXT, "align": PP_ALIGN.CENTER},
])

stat_y = CTY + INFO_CARD_H + CARD_GAP
rect(s3, rx, stat_y, rw, STAT_CARD_H, fill=NAVY)
vtxt(s3, rx, stat_y + 0.2, rw, 0.7,
     "48,806", size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER, valign="b")
vtxt(s3, rx, stat_y + 0.95, rw, 0.55,
     "Titik Pengamatan", size=12, color=LIGHT_AQU,
     align=PP_ALIGN.CENTER, valign="t")
vtxt(s3, rx, stat_y + 1.3, rw, 0.4,
     "Gravitasi Airborne", size=11, color=LIGHT_AQU,
     align=PP_ALIGN.CENTER, valign="t")


# ──────────────────────────────────────────────
# SLIDE 4 — DATA SURVEY
# ──────────────────────────────────────────────
s4 = content_slide("Data Survei Gravitasi Airborne")

from pptx.oxml import parse_xml

def _cell_fill(cell, hex_str):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = parse_xml(
        f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:srgbClr val="{hex_str}"/></a:solidFill>'
    )
    for child in list(tcPr):
        if "Fill" in child.tag or "fill" in child.tag.lower():
            tcPr.remove(child)
    tcPr.insert(0, solidFill)

table_data = [
    ("Parameter",    "Nilai"),
    ("Jumlah titik", "48,806"),
    ("Titik darat",  "16,167  (33%)"),
    ("Titik laut",   "32,639  (67%)"),
    ("FAA range",    "-149.12  s/d  +240.63 mGal"),
    ("Elevasi",      "-3,952  s/d  +2,568 m"),
]
tbl = s4.shapes.add_table(
    len(table_data), 2,
    Inches(0.35), Inches(CTY + 0.05),
    Inches(5.1), Inches(3.8)
).table
tbl.columns[0].width = Inches(2.5)
tbl.columns[1].width = Inches(2.6)

for r_idx, (param, val) in enumerate(table_data):
    for c_idx, txt in enumerate((param, val)):
        cell = tbl.cell(r_idx, c_idx)
        cell.text = txt
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.name = "Calibri"
        run.font.size = Pt(12)
        if r_idx == 0:
            run.font.bold = True
            run.font.color.rgb = WHITE
            _cell_fill(cell, "1A2A5E")
        elif r_idx % 2 == 0:
            run.font.color.rgb = DARK_TEXT
            _cell_fill(cell, "EFF6FF")
        else:
            run.font.color.rgb = DARK_TEXT

# Right panel — data columns
col_defs = [
    ("Bujur (Longitude)",  "Posisi bujur (°)"),
    ("Lintang (Latitude)", "Posisi lintang (°)"),
    ("FAA",                "Free Air Anomaly (mGal)"),
    ("Elevation_m",        "Ketinggian (m)"),
    ("UTM_X / UTM_Y",      "Koordinat UTM Zone 51S"),
]
rx = 5.75
rect(s4, rx, CTY, 3.9, 0.42, fill=NAVY)
vtxt(s4, rx, CTY, 3.9, 0.42, "Kolom Data",
     size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, valign="ctr")

for i, (name, desc) in enumerate(col_defs):
    yy = CTY + 0.42 + i * 0.62
    bg_col = WHITE if i % 2 == 0 else LIGHT_BLU
    rect(s4, rx, yy, 3.9, 0.56, fill=bg_col, line_color=BORDER, line_pt=0.5)
    rect(s4, rx, yy, 0.07, 0.56, fill=TEAL)
    multi_para(s4, rx + 0.18, yy + 0.07, 3.6, 0.46, [
        {"text": name, "size": 12, "bold": True, "color": NAVY, "space_after": 2},
        {"text": desc, "size": 10, "color": MUTED},
    ])

vtxt(s4, 0.35, SH - 0.3, 9.3, 0.22,
     "Sumber data: gravity_filled.csv",
     size=9, italic=True, color=MUTED, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────
# SLIDE 5 — METODOLOGI  (flowchart)
# ──────────────────────────────────────────────
s5 = content_slide("Metodologi Pengolahan")

steps = [
    "Data Loading & QC",
    "Persiapan DEM (BATNAS)",
    "Koreksi Bouguer Sederhana\n(harmonica.bouguer_correction)",
    "Koreksi Terrain\n(Prism Forward Modeling)",
    "Gridding\n(Equivalent Sources)",
    "Peta Anomali Gravitasi",
]
bw, bh = 3.4, 0.54
bx = 0.45
arr    = 0.22
step_h = bh + arr
start_y = CTY + 0.05

for i, step in enumerate(steps):
    y = start_y + i * step_h
    is_last = (i == len(steps) - 1)
    fill_c = NAVY if is_last else LIGHT_BLU
    line_c = NAVY if is_last else TEAL
    txt_c  = WHITE if is_last else DARK_TEXT
    rect(s5, bx, y, bw, bh, fill=fill_c, line_color=line_c, line_pt=1.5)
    rect(s5, bx, y, 0.07, bh, fill=TEAL if not is_last else TEAL_MID)
    vtxt(s5, bx + 0.12, y, bw - 0.15, bh, step,
         size=11, bold=is_last, color=txt_c, valign="ctr")
    if not is_last:
        ay = y + bh
        rect(s5, bx + bw/2 - 0.02, ay, 0.04, arr * 0.65, fill=TEAL)
        vtxt(s5, bx + bw/2 - 0.15, ay + arr * 0.55, 0.3, 0.18,
             "▼", size=9, color=TEAL, align=PP_ALIGN.CENTER)

# Right panel — software
rx2 = 4.1
rw2 = 5.6
rect(s5, rx2, CTY, rw2, 0.38, fill=NAVY)
vtxt(s5, rx2, CTY, rw2, 0.38, "Software & Referensi",
     size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, valign="ctr")

software = [
    ("Python 3.12",          "Platform komputasi utama"),
    ("harmonica 0.7.0",      "Koreksi gravitasi & terrain"),
    ("verde 1.9.0",          "Equivalent Sources & gridding"),
    ("boule",                "Normal gravity reference (WGS84)"),
    ("matplotlib + cartopy", "Visualisasi dan peta"),
]
for i, (lib, desc) in enumerate(software):
    yy = CTY + 0.38 + i * 0.58
    bg = WHITE if i % 2 == 0 else RGBColor(0xF0, 0xF9, 0xFF)
    rect(s5, rx2 + 0.1, yy, rw2 - 0.2, 0.52, fill=bg, line_color=BORDER, line_pt=0.5)
    multi_para(s5, rx2 + 0.25, yy + 0.06, rw2 - 0.45, 0.44, [
        {"text": lib,  "size": 12, "bold": True, "color": NAVY, "space_after": 2},
        {"text": desc, "size": 10, "color": MUTED},
    ])

vtxt(s5, rx2 + 0.1, CTY + 0.38 + 5 * 0.58, rw2 - 0.2, 0.5,
     "Referensi: Geosoft Oasis Montaj workflow",
     size=10, italic=True, color=MUTED, align=PP_ALIGN.CENTER, valign="ctr")


# ──────────────────────────────────────────────
# SLIDES 6–10  (image + key points)
# Wide images (ratio > 1.7): full-width top image, 2-col bullets below
# Square/tall images (ratio ≤ 1.7): left image + right bullets panel
# ──────────────────────────────────────────────
def img_bullet_slide(title, img_file, caption, bullets):
    sl = content_slide(title)
    img_path = os.path.join(OUTPUT_DIR, img_file)
    _, _, ratio = img_ratio(img_path)

    if ratio > 1.7:
        # ── Wide layout ──────────────────────────────────────
        # Image fills full available width; height capped at 3.5"
        # leaving ~1.0" for 2-column bullets below
        max_w_img = SW - 0.7    # 9.3"
        max_h_img = 3.5
        bot = add_img(sl, img_path, 0.35, CTY, max_w_img, max_h_img)

        vtxt(sl, 0.35, bot + 0.04, max_w_img, 0.26, caption,
             size=9, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

        bul_y = bot + 0.34
        bul_h = SH - bul_y - 0.15
        n_left = (len(bullets) + 1) // 2
        left_rows  = [{"text": b, "size": 11, "color": DARK_TEXT, "space_after": 4}
                      for b in bullets[:n_left]]
        right_rows = [{"text": b, "size": 11, "color": DARK_TEXT, "space_after": 4}
                      for b in bullets[n_left:]]
        multi_para(sl, 0.45, bul_y, 4.4, bul_h, left_rows)
        multi_para(sl, 5.1,  bul_y, 4.4, bul_h, right_rows)

    else:
        # ── Left image + right key-points panel ──────────────
        max_w_img = 6.0
        max_h_img = SH - CTY - 0.35
        bot = add_img(sl, img_path, 0.25, CTY, max_w_img, max_h_img)

        vtxt(sl, 0.25, bot + 0.04, max_w_img, 0.27, caption,
             size=8.5, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

        rx = 6.45
        rw = SW - rx - MARG
        ry = CTY

        rect(sl, rx, ry, rw, 0.36, fill=NAVY)
        vtxt(sl, rx, ry, rw, 0.36, "Key Points",
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, valign="ctr")
        rect(sl, rx, ry + 0.36, rw, SH - ry - 0.36 - 0.18,
             fill=WHITE, line_color=BORDER, line_pt=0.75)
        rows = [{"text": b, "size": 11, "color": DARK_TEXT, "space_after": 7}
                for b in bullets]
        multi_para(sl, rx + 0.15, ry + 0.5, rw - 0.22,
                   SH - ry - 0.5 - 0.2, rows)
    return sl


img_bullet_slide(
    "Step 1 — Quality Control Data",
    "01_qc_report.png",
    "QC Report: Distribusi FAA dan Elevasi",
    [
        "Pemeriksaan statistik pada data mentah gravitasi",
        "Deteksi outlier dengan metode IQR (Interquartile Range)",
        "Threshold: 3 x IQR",
        "Tidak ditemukan outlier di luar threshold",
        "48,806 titik data lolos quality control",
    ],
)

img_bullet_slide(
    "Step 2 — Persiapan DEM (BATNAS)",
    "02_dem_preview.png",
    "DEM dan distribusi titik pengamatan gravitasi",
    [
        "Sumber: BATNAS v1.5 (BIG Indonesia)",
        "Resolusi: ~6 arc-second (~185 m)",
        "2 tile digabungkan, dipotong sesuai area survei",
        "Diproyeksikan ke UTM Zone 51S (EPSG:32751)",
        "Titik pengamatan gravitasi: titik merah pada peta",
    ],
)

img_bullet_slide(
    "Step 3 — Koreksi Bouguer Sederhana",
    "03_bouguer_anomaly.png",
    "Perbandingan FAA dan Simple Bouguer Anomaly (SBA)",
    [
        "Menghilangkan efek massa topografi dari FAA",
        "Densitas kerak: 2,670 kg/m3  |  Densitas air: 1,040 kg/m3",
        "Rumus: SBA = FAA - Bouguer Correction",
        "harmonica.bouguer_correction() — darat & laut otomatis",
        "Hasil: SBA range  -79.8  s/d  +185.5 mGal",
    ],
)

img_bullet_slide(
    "Step 4 — Koreksi Terrain (Prism Forward Modeling)",
    "04_terrain_correction.png",
    "Terrain correction, SBA, dan Complete Bouguer Anomaly",
    [
        "DEM dimodelkan sebagai layer prisma (prism_layer)",
        "Gravitasi topografi dihitung di tiap titik pengamatan",
        "Densitas: 2,670 kg/m3 (darat) | 1,630 kg/m3 (laut)",
        "CBA = FAA - efek terrain",
        "CBA lebih halus dari SBA (std 37.0 vs 38.8 mGal)",
    ],
)

img_bullet_slide(
    "Step 5 — Gridding (Equivalent Sources)",
    "05_gridded_maps.png",
    "Perbandingan data tersebar dan hasil gridding reguler",
    [
        "Metode: Equivalent Sources (harmonica + verde)",
        "Data titik tersebar > grid reguler 2 km x 2 km",
        "Ketinggian prediksi: 2,600 m (upward continuation)",
        "Damping: 10  |  Depth: 10,000 m",
        "Interpolasi menghormati fisika medan potensial",
    ],
)


# ──────────────────────────────────────────────
# SLIDES 11–13  (left image + right interpretation)
# ──────────────────────────────────────────────
def full_img_slide(title, img_file, caption, interp=None):
    sl = content_slide(title)
    img_path = os.path.join(OUTPUT_DIR, img_file)

    # Left image — 6.3" wide, full content height
    max_w_img = 6.3
    max_h_img = SH - CTY - 0.2
    bot = add_img(sl, img_path, 0.25, CTY, max_w_img, max_h_img)

    vtxt(sl, 0.25, bot + 0.04, max_w_img, 0.26, caption,
         size=9, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

    # Right panel
    rx  = 6.75
    rw  = SW - rx - 0.25
    ry  = CTY

    if interp:
        rect(sl, rx, ry, rw, 0.36, fill=NAVY)
        vtxt(sl, rx, ry, rw, 0.36, "Interpretasi",
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, valign="ctr")
        rect(sl, rx, ry + 0.36, rw, SH - ry - 0.36 - 0.18,
             fill=WHITE, line_color=BORDER, line_pt=0.75)
        rows = [{"text": b, "size": 12, "color": DARK_TEXT, "space_after": 14}
                for b in interp]
        multi_para(sl, rx + 0.15, ry + 0.5, rw - 0.22,
                   SH - ry - 0.5 - 0.2, rows)
    return sl


full_img_slide(
    "Peta Anomali Free Air (FAA)",
    "06_faa_anomaly_map.png",
    "Free Air Anomaly — Survei Gravitasi Airborne Sulawesi Tengah",
    [
        "Anomali positif tinggi (merah): area daratan dengan topografi tinggi",
        "Anomali negatif (biru): cekungan laut dalam",
        "Korelasi kuat dengan topografi — sesuai harapan untuk FAA",
    ],
)

full_img_slide(
    "Peta Complete Bouguer Anomaly (CBA)",
    "06_cba_anomaly_map.png",
    "Complete Bouguer Anomaly — variasi densitas bawah permukaan",
    [
        "Efek topografi telah dihilangkan sepenuhnya",
        "Anomali mencerminkan variasi densitas di bawah permukaan",
        "Berguna untuk identifikasi: sesar, intrusi, cekungan",
    ],
)

full_img_slide(
    "Hasil Keseluruhan",
    "06_combined_report.png",
    "(a) Topografi/Batimetri    (b) FAA    (c) CBA    (d) Anomali Residual",
    [
        "Panel (d): Anomali residual setelah menghilangkan tren regional"
        " — mengungkap fitur geologis lokal",
    ],
)


# ──────────────────────────────────────────────
# SLIDE 14 — CONCLUSION
# ──────────────────────────────────────────────
s14 = prs.slides.add_slide(BLANK)
set_bg(s14, NAVY)
rect(s14, 0, 0,        SW, 0.18, fill=TEAL)
rect(s14, 0, SH - 0.18, SW, 0.18, fill=TEAL)

vtxt(s14, 0.3, 0.22, SW - 0.6, 0.7,
     "Kesimpulan & Rencana Selanjutnya",
     size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER, valign="ctr")
rect(s14, 0.5, 1.0, 9.0, 0.04, fill=TEAL)

# Left: Kesimpulan
rect(s14, 0.3, 1.1, 4.5, 0.38, fill=TEAL)
vtxt(s14, 0.3, 1.1, 4.5, 0.38, "KESIMPULAN",
     size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, valign="ctr")
conclusions = [
    "Data gravitasi airborne berhasil diolah (48,806 titik)",
    "Peta FAA dan CBA berhasil dihasilkan",
    "Koreksi terrain via prism forward modeling",
    "Gridding menggunakan Equivalent Sources (2 km)",
]
multi_para(s14, 0.35, 1.58, 4.35, SH - 1.58 - 0.35,
           [{"text": "✓  " + c, "size": 13, "color": WHITE, "space_after": 20}
            for c in conclusions])

# Right: Rencana
rect(s14, 5.2, 1.1, 4.5, 0.38, fill=TEAL_MID)
vtxt(s14, 5.2, 1.1, 4.5, 0.38, "RENCANA SELANJUTNYA",
     size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER, valign="ctr")
next_steps = [
    "Pemisahan Regional-Residual (filtering frekuensi)",
    "Analisis derivatif (gradient horizontal, vertikal, tilt)",
    "Euler Deconvolution (estimasi kedalaman sumber)",
    "Analisis spektrum daya (estimasi kedalaman)",
    "Profil penampang 2D",
    "Interpretasi geologi",
]
multi_para(s14, 5.25, 1.58, 4.35, SH - 1.58 - 0.35,
           [{"text": "->  " + s, "size": 13, "color": WHITE, "space_after": 13}
            for s in next_steps])

# Bottom "Thank you" note to fill lower whitespace
vtxt(s14, 0.5, SH - 0.72, SW - 1.0, 0.48,
     "Terima Kasih",
     size=18, bold=True, color=TEAL_MID,
     align=PP_ALIGN.CENTER, valign="ctr")


# ──────────────────────────────────────────────
# SAVE
# ──────────────────────────────────────────────
prs.save(OUTPUT_FILE)
print(f"[OK] Saved: {OUTPUT_FILE}")
print(f"     Slides: {len(prs.slides)}")
print(f"     Size:   {os.path.getsize(OUTPUT_FILE) / 1024:.0f} KB")
